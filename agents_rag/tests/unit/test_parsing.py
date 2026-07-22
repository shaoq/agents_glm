"""parsing + cleaning 测试（docling PDF 解析留端到端验证）。"""

from __future__ import annotations

import pytest

from agents_rag.cleaning.normalizer import Normalizer, clean_text
from agents_rag.models import Block, BlockType, DocType, Document, Section
from agents_rag.parsing.base import Parser
from agents_rag.parsing.html_parser import HTMLParser
from agents_rag.parsing.markdown_parser import MarkdownParser, md_to_sections
from agents_rag.parsing.office_parser import DocxParser
from agents_rag.parsing.quality import assess_quality
from agents_rag.parsing.router import ParserRouter


# —— markdown ——
def test_markdown_headings_levels():
    sections = md_to_sections("# H1\na\n\n## H2\nb\n")
    assert sections[0].heading == "H1" and sections[0].level == 1
    assert sections[0].children[0].heading == "H2"
    assert sections[0].children[0].level == 2


def test_markdown_table_block():
    md = "# T\n| a | b |\n|---|---|\n| 1 | 2 |\n"
    doc = Document(
        doc_id="", source="s", doc_type=DocType.MARKDOWN, sections=tuple(md_to_sections(md))
    )
    tables = [b for b in doc.iter_blocks() if b.type is BlockType.TABLE]
    assert tables and "a" in tables[0].text


def test_markdown_multilevel_section_path():
    from agents_rag.chunking.structural import split_parents

    md = (
        "# H1\n" + "正文一" * 30 + "\n\n"
        "## H2\n" + "正文二" * 30 + "\n\n"
        "### H3\n" + "正文三" * 30 + "\n"
    )
    doc = Document(
        doc_id="d", source="s", doc_type=DocType.MARKDOWN, sections=tuple(md_to_sections(md))
    )
    parents = split_parents(doc, parent_max_size=100000)
    paths = {p.section_path for p in parents}
    assert any("H1" in p for p in paths)
    assert any("H2" in p for p in paths)
    assert any("H3" in p for p in paths)
    assert "H1 > H2 > H3" in paths  # 完整层级路径


def test_markdown_parser_document(tmp_path):
    f = tmp_path / "a.md"
    f.write_text("# H\n正文\n", encoding="utf-8")
    doc = MarkdownParser().parse(f)
    assert doc.doc_type is DocType.MARKDOWN
    assert any(b.text == "正文" for b in doc.iter_blocks())


# —— html ——
def test_html_parser_extracts_body(tmp_path):
    f = tmp_path / "a.html"
    f.write_text(
        "<html><body><nav>导航</nav><article><h1>T</h1><p>正文内容</p></article></body></html>",
        encoding="utf-8",
    )
    doc = HTMLParser().parse(f)
    text = " ".join(b.text for b in doc.iter_blocks())
    assert "正文内容" in text


# —— office docx ——
def _make_docx(path) -> None:
    import docx

    d = docx.Document()
    d.add_paragraph("第一段")
    d.add_paragraph("第二段")
    tbl = d.add_table(rows=2, cols=2)
    tbl.rows[0].cells[0].text = "k1"
    tbl.rows[0].cells[1].text = "k2"
    tbl.rows[1].cells[0].text = "v1"
    tbl.rows[1].cells[1].text = "v2"
    d.save(str(path))


def test_docx_parser_paragraphs_and_table(tmp_path):
    p = tmp_path / "a.docx"
    _make_docx(p)
    doc = DocxParser().parse(p)
    texts = [b.text for b in doc.iter_blocks()]
    assert "第一段" in texts
    tables = [b for b in doc.iter_blocks() if b.type is BlockType.TABLE]
    assert tables and "k1" in tables[0].text


# —— quality ——
def test_assess_quality_poor_for_scan_like():
    doc = Document(
        doc_id="",
        source="s",
        doc_type=DocType.PDF,
        sections=(Section(blocks=(Block(text="ab", page=1), Block(text="cd", page=2))),),
    )
    assert assess_quality(doc).is_poor()  # 每页字符极少


def test_poor_reason_diagnosis():
    from agents_rag.models import QualityReport

    assert QualityReport(chars_per_page=5).poor_reason() == "scan"  # 极低
    assert QualityReport(chars_per_page=200, garbage_ratio=0.5).poor_reason() == "layout"
    assert QualityReport(chars_per_page=200, garbage_ratio=0.1).poor_reason() is None
    # scan 阈值严：chars=30 是 is_poor 但非 poor_reason（不降级）
    assert QualityReport(chars_per_page=30).is_poor()
    assert QualityReport(chars_per_page=30).poor_reason() is None


# —— normalizer ——
def test_normalizer_preserves_page_metadata():
    b = Block(type=BlockType.PARAGRAPH, text="你好　　世界", page=3)  # 全角双空格
    doc = Document(
        doc_id="d",
        source="s",
        doc_type=DocType.MARKDOWN,
        sections=(Section(blocks=(b,)),),
    )
    out = Normalizer().normalize(doc)
    ob = out.sections[0].blocks[0]
    assert ob.page == 3  # 元数据保留
    assert "　　" not in ob.text  # 全角空格被压缩


def test_clean_text_keeps_symbols():
    assert "#" in clean_text("a # b")
    assert "1.0" in clean_text("型号 A-1.0 编码")


# —— router（诊断式）——
class _BoomParser(Parser):
    supported_types = (DocType.TXT,)

    def parse(self, path):  # type: ignore[override]
        raise RuntimeError("boom")


class _FakeParser(Parser):
    """产出指定 Document，用于控制 poor_reason。"""

    def __init__(self, doc: Document):
        self._doc = doc

    def parse(self, path):  # type: ignore[override]
        return self._doc


def _txt_doc(text: str, page: int | None = None) -> Document:
    return Document(
        doc_id="", source="s", doc_type=DocType.TXT,
        sections=(Section(blocks=(Block(text=text, page=page),)),),
    )


def test_router_primary_ok(tmp_path):
    f = tmp_path / "a.txt"; f.write_text("x")
    ok = _txt_doc("正常文本内容足够长不至于被判低质量" * 5)
    router = ParserRouter(primary={DocType.TXT: _FakeParser(ok)})
    assert router.parse(f) is ok  # 达标 → 返回主 doc


def test_router_primary_fail_returns_none(tmp_path):
    f = tmp_path / "a.txt"; f.write_text("x")
    router = ParserRouter(primary={DocType.TXT: _BoomParser()})
    assert router.parse(f) is None  # 主 raise → None


def test_router_no_parser_for_type(tmp_path):
    f = tmp_path / "a.xyz"; f.write_text("x")
    router = ParserRouter(primary={})
    assert router.parse(f) is None


def test_router_poor_scan_falls_back(tmp_path):
    f = tmp_path / "a.txt"; f.write_text("x")
    scan_doc = _txt_doc("ab", page=1)  # chars/page=2 < 10 → scan
    ok_doc = _txt_doc("OCR 识别出的足够长正文内容" * 5)
    router = ParserRouter(
        primary={DocType.TXT: _FakeParser(scan_doc)},
        fallbacks={DocType.TXT: {"scan": _FakeParser(ok_doc)}},
    )
    assert router.parse(f) is ok_doc  # scan → OCR fallback


def test_router_poor_no_fallback_skip(tmp_path):
    f = tmp_path / "a.txt"; f.write_text("x")
    scan_doc = _txt_doc("ab", page=1)  # scan poor
    router = ParserRouter(primary={DocType.TXT: _FakeParser(scan_doc)})
    assert router.parse(f) is None  # poor + 无兜底 → None


def test_router_poor_fallback_still_poor_skip(tmp_path):
    f = tmp_path / "a.txt"; f.write_text("x")
    scan_doc = _txt_doc("ab", page=1)
    bad_fb = _txt_doc("xy", page=1)  # 兜底也 poor
    router = ParserRouter(
        primary={DocType.TXT: _FakeParser(scan_doc)},
        fallbacks={DocType.TXT: {"scan": _FakeParser(bad_fb)}},
    )
    assert router.parse(f) is None  # 兜底仍 poor → None


def test_xlsx_parser(tmp_path):
    from openpyxl import Workbook

    from agents_rag.parsing.office_parser import XlsxParser

    wb = Workbook()
    ws = wb.active
    ws.append(["型号", "参数"])
    ws.append(["GLM-4.5", "2048维"])
    p = tmp_path / "a.xlsx"
    wb.save(str(p))
    doc = XlsxParser().parse(p)
    tables = [b for b in doc.iter_blocks() if b.type is BlockType.TABLE]
    assert tables and "GLM-4.5" in tables[0].text


def test_pptx_parser(tmp_path):
    from pptx import Presentation

    from agents_rag.parsing.office_parser import PptxParser

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(0, 0, 300, 100)
    tb.text_frame.text = "智谱GLM模型说明文本内容"
    p = tmp_path / "a.pptx"
    prs.save(str(p))
    doc = PptxParser().parse(p)
    assert "GLM" in " ".join(b.text for b in doc.iter_blocks())


def test_router_with_defaults_disables_pdf():
    r = ParserRouter.with_defaults(enable_pdf=False)
    assert DocType.MARKDOWN in r._primary
    assert DocType.PDF not in r._primary  # docling 未启用
