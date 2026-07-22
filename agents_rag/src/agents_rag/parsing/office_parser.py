"""Office 解析：docx / xlsx / pptx → sections + 表格 block。笔记 §3.3。"""

from __future__ import annotations

from pathlib import Path

from agents_rag.models import Block, BlockType, Document, DocType, Section
from agents_rag.parsing.base import Parser


def rows_to_markdown(rows: list[list[str]]) -> str:
    """行列矩阵 → markdown 表（首行为表头）。"""
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    norm = [(r + [""] * width)[:width] for r in rows]
    header = norm[0]
    lines = ["| " + " | ".join(header) + " |"]
    lines.append("| " + " | ".join("---" for _ in header) + " |")
    for r in norm[1:]:
        lines.append("| " + " | ".join(r) + " |")
    return "\n".join(lines)


def _doc_from_blocks(source: str | Path, doctype: DocType, blocks: list[Block]) -> Document:
    blocks = tuple(b for b in blocks if b.text)
    return Document(
        doc_id="",
        source=str(source),
        doc_type=doctype,
        sections=(Section(heading=None, level=1, blocks=blocks),),
    )


class DocxParser(Parser):
    supported_types = (DocType.DOCX,)

    def parse(self, path: str | Path) -> Document:
        import docx

        d = docx.Document(str(path))
        blocks: list[Block] = []
        for p in d.paragraphs:
            t = p.text.strip()
            if t:
                blocks.append(Block(type=BlockType.PARAGRAPH, text=t))
        for tbl in d.tables:
            rows = [[c.text.strip() for c in row.cells] for row in tbl.rows]
            blocks.append(
                Block(
                    type=BlockType.TABLE,
                    text=rows_to_markdown(rows),
                    table_data={"rows": rows},
                )
            )
        return _doc_from_blocks(path, DocType.DOCX, blocks)


class XlsxParser(Parser):
    supported_types = (DocType.XLSX,)

    def parse(self, path: str | Path) -> Document:
        from openpyxl import load_workbook

        wb = load_workbook(str(path), read_only=True, data_only=True)
        blocks: list[Block] = []
        for ws in wb.worksheets:
            rows = [
                [("" if c is None else str(c)) for c in row]
                for row in ws.iter_rows(values_only=True)
                if any(c is not None for c in row)
            ]
            if rows:
                blocks.append(
                    Block(
                        type=BlockType.TABLE,
                        text=rows_to_markdown(rows),
                        table_data={"rows": rows},
                    )
                )
        wb.close()
        return _doc_from_blocks(path, DocType.XLSX, blocks)


class PptxParser(Parser):
    supported_types = (DocType.PPTX,)

    def parse(self, path: str | Path) -> Document:
        from pptx import Presentation

        prs = Presentation(str(path))
        blocks: list[Block] = []
        for slide in prs.slides:
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        texts.append(t)
            if texts:
                blocks.append(Block(type=BlockType.PARAGRAPH, text="\n".join(texts)))
        return _doc_from_blocks(path, DocType.PPTX, blocks)
